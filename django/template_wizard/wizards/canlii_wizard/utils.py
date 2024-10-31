import asyncio
import concurrent.futures
import io
import math
import os
import re
import tempfile
from datetime import datetime as dt
from uuid import uuid4

from django.conf import settings

import docx
import markdown
import tiktoken
from bs4 import BeautifulSoup
from docxtpl import DocxTemplate
from langchain.chains import ReduceDocumentsChain
from langchain.chains.combine_documents.stuff import StuffDocumentsChain
from langchain.chains.llm import LLMChain
from langchain.prompts import PromptTemplate
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import AzureChatOpenAI
from openai import AzureOpenAI
from pptx import Presentation

client = AzureOpenAI(
    api_key=settings.AZURE_OPENAI_KEY,
    api_version=settings.AZURE_OPENAI_VERSION,
    azure_endpoint=settings.AZURE_OPENAI_ENDPOINT,
)

model_str = settings.DEFAULT_CHAT_MODEL

system_prompt = {
    "en": "You are a legal professional who is tasked with reviewing a case and writing legal case reports. You're going to extract and write information from the following content and will only use the facts provided in the content I provide you. When providing responses, please use paragraphs instead of bullet points. Do not include question marks or exclamation marks, and ensure a constant professional style of writing throughout all your responses.",
    "fr": "Vous êtes un professionnel du droit chargé d'examiner un dossier et de rédiger des rapports juridiques. Vous allez extraire et rédiger des informations à partir du contenu suivant et n'utiliserez que les faits fournis dans le contenu que je vous fournis. Lorsque vous fournissez des réponses, veuillez utiliser des paragraphes au lieu de points de puce. N'incluez pas de points d'interrogation ou d'exclamation, et assurez-vous d'avoir un style d'écriture professionnel constant dans toutes vos réponses.",
}

# use this for french months
months_in_french = {
    "january": "janvier",
    "february": "février",
    "march": "mars",
    "april": "avril",
    "may": "mai",
    "june": "juin",
    "july": "juillet",
    "august": "août",
    "september": "septembre",
    "october": "octobre",
    "november": "novembre",
    "december": "décembre",
}
acronym_dict = {
    "IRB": "Immigration and Refugee Board",
    "RPD": "Refugee Protection Division of the IRB",
    "RAD": "Refugee Appeal Division of the IRB",
    "ID": "Immigration Division of the IRB",
    "IAD": "Immigration Appeal Division of the IRB",
    "IRCC": "Immigration, Refugees and Citizenship Canada",
    "IRPA": "Immigration and Refugee Protection Act",
    "IRP Regulations": "Immigration and Refugee Protection Regulations",
    "FOSS": "Field Operational Support System",
    "CAIPS": "Computer Assisted Immigration Processing System",
    "GCMS": "Global Case Management System",
    "MCI": "Minister of Citizenship and Immigration",
    "CIC": "Citizenship and Immigration Canada",
    "PRRA": "Pre-Removal Risk Assessment",
    "H&amp;C": "Humanitarian and Compassionate",
    "BIOC": "Best interests of the child",
    "DCO": "Designated Country of Origin",
    "MPSEP": "Minister of Public Safety and Emergency Preparedness",
    "CBSA": "Canada Border Services Agency",
}


def map_reduce(text, language, model_str):

    length_prompt_en = "Write down a long document with large paragraphs containing multiple sentences that detail all important details of the text (in English). Simply rewrite; do not say 'This document is about...' etc, this not just a summary but a detailed rewording of the information presented. There is no length limit - be as detailed as possible. However, **do not extrapolate** on the text. The document that you will write must be factual and not introduce any new ideas. Remember, do not list information in bullet points, make sure the response is in the form of a long documents with multiple paragraphs."

    length_prompt_fr = "Écrivez un long document avec de grands paragraphes contenant plusieurs phrases qui détaillent tous les détails importants du texte (en anglais). Il suffit de le réécrire; ne dites pas 'Ce document parle de...' etc, ce n'est pas seulement un résumé mais une reformulation détaillée des informations présentées. Il n'y a pas de limite de longueur - soyez aussi détaillé que possible. Cependant, n'extrapolez pas sur le texte. Le document que vous écrirez doit être factuel et ne doit introduire aucune nouvelle idée. N'oubliez pas, ne répertoriez pas les informations sous forme de points, assurez-vous que la réponse est sous forme d'un long document avec plusieurs paragraphes."

    llm = AzureChatOpenAI(
        azure_endpoint=settings.AZURE_OPENAI_ENDPOINT,
        azure_deployment=model_str,
        model=model_str,
        api_version=settings.AZURE_OPENAI_VERSION,
        api_key=settings.AZURE_OPENAI_KEY,
        temperature=0.3,
    )

    reduce_template_en = (
        "The following are parts of a larger document:\n\n"
        "{docs}\n\n"
        f"{length_prompt_en}. Return it in markdown format.\n"
        "Helpful Answer:"
    )
    reduce_template_fr = (
        "Les résumés suivants ont été générés à partir d'un document:\n\n"
        "{docs}\n\n"
        f"Prenez-les et réécrivez-le en un résumé final. {length_prompt_fr}. Renvoyez-le au format markdown.\n"
        "Réponse utile:"
    )
    reduce_template = reduce_template_fr if language == "fr" else reduce_template_en
    reduce_prompt = PromptTemplate.from_template(reduce_template)

    # Run chain
    reduce_chain = LLMChain(llm=llm, prompt=reduce_prompt)

    # Takes a list of documents, combines them into a single string, and passes this to an LLMChain
    combine_documents_chain = StuffDocumentsChain(
        llm_chain=reduce_chain, document_variable_name="docs"
    )

    # Combines and iteratively reduces the mapped documents
    reduce_documents_chain = ReduceDocumentsChain(
        # This is final chain that is called.
        combine_documents_chain=combine_documents_chain,
        # If documents exceed context for `StuffDocumentsChain`
        collapse_documents_chain=combine_documents_chain,
        # The maximum number of tokens to group documents into.
        token_max=15000,
    )

    # Split the text into chunks
    text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=12000, chunk_overlap=100
    )

    docs = [Document(page_content=text, metadata={"source": "userinput"})]

    split_docs = text_splitter.split_documents(docs)

    async def run_map_reduce():

        # Run the chain asynchronously and return the result when it's done
        return await reduce_documents_chain.arun(split_docs)

    response = asyncio.run(run_map_reduce())

    return response


def trim_to_tokens(text, max_tokens=15000):
    # Initialize the tokenizer with the specified encoding
    enc = tiktoken.get_encoding("cl100k_base")

    # Encode the input text, getting token IDs
    token_ids = enc.encode(text)

    # Trim token IDs to get the beginning and end parts as specified
    tokens_beginning = (
        token_ids[:max_tokens] if len(token_ids) > max_tokens else token_ids
    )
    tokens_end = token_ids[-max_tokens:] if len(token_ids) > max_tokens else []

    # Decode the trimmed token IDs back into text
    trimmed_beginning = enc.decode(tokens_beginning)
    trimmed_end = enc.decode(tokens_end) if tokens_end else ""

    return trimmed_beginning, trimmed_end


def extract_text_from_pdf(pdf_file):
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(pdf_file)
    text = ""
    for i, page in enumerate(pdf):
        text_page = page.get_textpage()
        text += text_page.get_text_range() + "\n"
        text_page.close()
    pdf.close()
    return text


def convert_to_text(file, content_type=None):
    text = ""
    if content_type is None:
        content_type = file.content_type
    if (
        content_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        doc = docx.Document(file)
        fullText = []
        for para in doc.paragraphs:
            fullText.append(para.text)
        text = "\n".join(fullText)
        text = text.replace(".", ". ")

    elif content_type == "application/pdf":
        text = extract_text_from_pdf(file)

    elif content_type == "text/plain":
        text = file.read().decode("utf-8")

    # Extract text from PowerPoint file
    elif (
        content_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "

    elif content_type == "text/html":
        file_content = file.read()
        soup = BeautifulSoup(file_content, "html.parser")

        # kill all script and style elements
        for script in soup(["script", "style"]):
            script.decompose()  # rip it out

        # get text
        text = soup.get_text()
        text = text.replace(".", ". ")

        # break into lines and remove leading and trailing space on each
        lines = (line.strip() for line in text.splitlines())
        # break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # drop blank lines
        text = "\n".join(chunk for chunk in chunks if chunk)

    return text


def extract_sections(input_string, sections=["Facts", "Findings"]):
    data = {section: None for section in sections}

    for i, section in enumerate(sections):
        if f"{section}:" in input_string:
            start = input_string.index(f"{section}:") + len(f"{section}:")
            end = len(input_string)
            if i < len(sections) - 1:  # If not the last section
                next_section = sections[i + 1]
                if f"{next_section}:" in input_string:
                    end = input_string.index(f"{next_section}:")
            data[section] = input_string[start:end].strip()

    return data


def extract_using_generative(extraction_word, content, model_link, language):

    question = {
        "en": f"""Extract the {extraction_word} from this case. Only write down the {extraction_word}, don't go into more details. Don't write down "The {extraction_word} is ...", simply writing down the {extraction_word} itself is enough.""",
        "fr": f"""Extrait {extraction_word} de ce cas. Écris seulement {extraction_word}, ne donne pas plus de détails. N'écris pas "{extraction_word} est ..." simplement écrire {extraction_word} sans phrase de contexte est suffisant.""",
    }
    extraction = client.chat.completions.create(
        model=model_link,  # model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": system_prompt[language]},
            {"role": "user", "content": "### Content: " + content},
            {"role": "user", "content": "### Question: " + question[language]},
        ],
        temperature=0.1,  # set this to a low value, such as 0.2, to make the output more focused and deterministic. This will reduce the randomness in the generated response
        top_p=1.0,  # set this to a high value, such as 0.9 or 1, so that the model only considers the tokens with the highest probability mass. This will make the output more predictable and less likely to generate unexpected responses.
        frequency_penalty=0.0,  # set this to a high positive value, such as 2.0, to penalize the model for repeating the same line verbatim. This will encourage the model to generate more diverse responses.
        presence_penalty=1.8,  # set this to a high positive value, such as 2.0, to penalize the model for introducing new tokens that do not appear in the input text. This will encourage the model to generate responses that are more closely related to the input text.
    )

    return extraction.choices[0].message.content, extraction.usage.total_tokens


def extract_court_judgment(content, model_str, language):

    question = {
        "en": """Find the complete court's judgment found under "THIS COURT’S JUDGMENT is that:". Make sure each individual court judgment is listed and numbered, in the form of 1. First court judgment 2. Second court judgment, and so forth. Your reply should simply be the complete list of the court judgment, don't include anything else.""",
        "fr": """Trouvez le jugement complet du tribunal sous "LE JUGEMENT DE CE TRIBUNAL est que:". Assurez-vous que chaque jugement individuel du tribunal est répertorié et numéroté, sous la forme de 1. Premier jugement du tribunal 2. Deuxième jugement du tribunal, et ainsi de suite. Votre réponse devrait simplement être la liste du jugement complet du tribunal, ne comprenant rien d'autre.""",
    }
    extraction = client.chat.completions.create(
        model=model_str,
        messages=[
            {"role": "system", "content": system_prompt[language]},
            {"role": "user", "content": "### Content: " + content},
            {"role": "user", "content": "### Question: " + question[language]},
        ],
        temperature=0.1,  # set this to a low value, such as 0.2, to make the output more focused and deterministic. This will reduce the randomness in the generated response
        frequency_penalty=0.0,  # set this to a high positive value, such as 2.0, to penalize the model for repeating the same line verbatim. This will encourage the model to generate more diverse responses.
        presence_penalty=1.8,  # set this to a high positive value, such as 2.0, to penalize the model for introducing new tokens that do not appear in the input text. This will encourage the model to generate responses that are more closely related to the input text.
    )

    return extraction.choices[0].message.content, extraction.usage.total_tokens


def replace_with_acronyms(text, acronym_dict):
    for acronym, full_version in acronym_dict.items():
        # Pattern to match the full version followed by an optional any text in parentheses
        regex_pattern = r"\b" + re.escape(full_version) + r"\b(?: \([^)]*\))?"
        # Replace the pattern with just the acronym
        text = re.sub(regex_pattern, acronym, text, flags=re.IGNORECASE)
    return text


def extract_meta_info(text_for_regex, text_for_llm, model_str, language):
    extra_cost = 0
    # Extract date
    date_match = re.search(r"Date:\s*(\d{4})((\s*(\d{1,2}))+)", text_for_regex)
    if date_match:
        matched_date = (
            "".join(date_match.groups())
            .replace("\n", "")
            .replace("\t", "")
            .replace(" ", "")
        )
        # Split based on number of digits
        year = int(matched_date[:4])
        month = int(matched_date[4:6])
        day = int(matched_date[6:8])
        date_obj = dt(year, month, day)

        if (
            language == "fr"
            and date_obj.strftime("%B").lower() in months_in_french.keys()
        ):
            # Get the month name in French
            month_name_french = months_in_french[date_obj.strftime("%B").lower()]

            # Format the date manually in French
            date = f"{day} {month_name_french.title()} {year}"  # Format like '31 Octobre 2023'
        else:

            # Format the date
            date = date_obj.strftime("%B %d, %Y")  # Format like 'October 23, 2023'

    else:
        extraction_word = "date" if language == "en" else "la date"
        date, date_cost = extract_using_generative(
            extraction_word, text_for_llm, model_str, language
        )
        extra_cost += date_cost

    # Extract docket number
    docket_match = re.search(
        r"(Docket|DOCKET):\s*([A-Z]+)\s*-\s*(\d+)\s*-\s*(\d+)", text_for_regex
    )
    if docket_match:
        docket_number = "-".join(
            docket_match.groups()[1:]
        )  # Excluding the first group as it's "Docket" or "DOCKET"
    else:
        extraction_word = (
            "docket number" if language == "en" else "le numéro de dossier"
        )
        docket_number, docket_number_cost = extract_using_generative(
            extraction_word, text_for_llm, model_str, language
        )
        extra_cost += docket_number_cost

    # Extract citation
    citation_match = re.search(
        r"Citation:\s*(\d{4})\s*([A-Z]+)\s*(\d+)", text_for_regex
    )
    if citation_match:
        citation = " ".join(citation_match.groups())
    else:
        extraction_word = "citation" if language == "en" else "la référence"
        citation, citation_cost = extract_using_generative(
            extraction_word, text_for_llm, model_str, language
        )
        extra_cost += citation_cost
    citation.replace("\n", " ").strip()

    # Extract the title
    question = {
        "en": "Find the title of this legal case. Write it in the format with the first entity before 'v.' and the second entity after. For example, Applicant v. Respondent. Only output the title itself, don't write anything else. Also make sure the reponse is in the format of a title with the appropirate capitalization.",
        "fr": "Trouvez le titre de cette affaire juridique. Écrivez-le dans le format avec la première entité avant 'v.' et la deuxième entité après. Par exemple, Demandeur v. Défendeur. Ne mentionnez que le titre lui-même, n'écrivez rien d'autre. Assurez-vous également que la réponse est dans le format d'un titre avec la capitalisation appropriée.",
    }

    title = client.chat.completions.create(
        model=model_str,
        messages=[
            {"role": "system", "content": system_prompt[language]},
            {"role": "user", "content": "### Content: " + text_for_llm},
            {"role": "user", "content": "### Question: " + question[language]},
        ],
        temperature=0.1,
        frequency_penalty=0.0,
        presence_penalty=1.8,
    )
    title = title.choices[0].message.content

    return date, docket_number, citation, title, extra_cost


def generate_immigration_case_report(text, language="en"):

    extra_cost = 0

    enc = tiktoken.get_encoding("cl100k_base")
    using_large_doc = False

    # Count how many characters are in the text
    token_count = len(enc.encode(text))

    # TODO: Change this to gpt-4 once we figure out how to deal with the gpt-4 delay errors
    token_count_max = 15000
    if token_count < token_count_max:
        content = text
    else:
        content = map_reduce(text, language, model_str)
        using_large_doc = True

    beginning_of_doc, end_of_doc = trim_to_tokens(text, max_tokens=token_count_max)

    # if we are using a large doc we'll just use the beginning of the doc to find what we want with the llm calls
    text_for_llm_meta_extraction = beginning_of_doc if using_large_doc else content
    date, docket_number, citation, title, extra_cost_meta = extract_meta_info(
        text, text_for_llm_meta_extraction, model_str, language
    )
    extra_cost += extra_cost_meta

    applicant = title.split("v.")[0].strip()
    respondent = title.split("v.")[1].strip()

    # Extract judge name
    judge_match = re.search(r"JUDGMENT AND REASONS:\s*([\w\s\.]+)DATED:", text)
    if judge_match:
        judge = judge_match.group(1).strip().title()
    else:
        judge_match = re.search(r"PRESENT:\s*The Honourable\s*([\w\s]+)BETWEEN:", text)
        if judge_match:
            judge = judge_match.group(1).strip().title()
        else:
            extraction_word = "judge" if language == "en" else "le juge"
            judge, judge_cost = extract_using_generative(
                extraction_word, text_for_llm_meta_extraction, model_str, language
            )
            extra_cost += judge_cost

    # if we are dealing with an english report lets try extracting the text itself first
    court_judgment_match = re.search(
        r"THIS COURT’S JUDGMENT is that\s*((.|\n)*)“[^”]+”\s*Judge", text
    )
    if court_judgment_match and language == "en":
        court_judgment = " ".join(court_judgment_match.groups())
    else:
        text_for_court_judgment_llm_call = end_of_doc if using_large_doc else content
        court_judgment, court_judgment_cost = extract_court_judgment(
            text_for_court_judgment_llm_call, model_str, language
        )

        extra_cost += court_judgment_cost

    def create_facts_and_court_sections():
        question = {
            "en": """Write down two section where the first section showcases the facts of the case, and the second section showcases the findings from the court. In both of these sections, make sure to mention in your answers important sections and subsections of relevant acts or regulations found in the text. For the first section, note that the facts are the “who, when, what, where, and why” of the case. Hence, describe the history of the applicant, the events that led to the case, the legal claims, and defenses of each party. As for the second section, summarize all the findings made specifically by the court during the trial, in addition to the result of the trial. If the findings did not come specifically from the court, do not mention it in the second section. Since all the findings mentioned in the second section came from the court, they should be written with the phrase "The Court" at the beginning of each sentence, for example "The Court found that...", "The Court noted..." or "The Court preferred...". Write the first section under "Facts:", and the second section under "Court:".""",
            "fr": """Écrivez deux sections où la première section présente les faits de l'affaire et la deuxième section présente les conclusions du tribunal. Dans ces deux sections, assurez-vous de mentionner dans vos réponses les sections et sous-sections importantes des lois ou règlements pertinents trouvés dans le texte. Pour la première section, notez que les faits sont le "qui, quand, quoi, où et pourquoi" de l'affaire. Par conséquent, décrivez l'historique du demandeur, les événements qui ont conduit à l'affaire, les revendications juridiques et les défenses de chaque partie. Quant à la deuxième section, résumez toutes les conclusions spécifiquement formulées par le tribunal lors du procès, ainsi que le résultat du procès. Si les conclusions ne proviennent pas spécifiquement du tribunal, ne les mentionnez pas dans la deuxième section. Étant donné que toutes les conclusions mentionnées dans la deuxième section proviennent du tribunal, elles doivent être rédigées avec la phrase "Le tribunal a constaté que", "Le tribunal a noté" ou "Le tribunal a préféré". Écrivez la première section sous "Faits:" et la deuxième section sous "Tribunal:".""",
        }

        facts = client.chat.completions.create(
            model=model_str,  # model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,  # set this to a low value, such as 0.2, to make the output more focused and deterministic. This will reduce the randomness in the generated response
            frequency_penalty=0.0,  # set this to a high positive value, such as 2.0, to penalize the model for repeating the same line verbatim. This will encourage the model to generate more diverse responses.
            presence_penalty=1.8,  # set this to a high positive value, such as 2.0, to penalize the model for introducing new tokens that do not appear in the input text. This will encourage the model to generate responses that are more closely related to the input text.
        )
        return facts

    def create_key_findings(findings):
        question = {
            "en": f"""Give a brief overview of the key finding from the text "{findings}" without going into details. Make sure your response has a length of 30 words or less.""",
            "fr": f"""Donnez un bref aperçu des principales conclusions du texte "{findings}" sans entrer dans les détails. Assurez-vous que votre réponse a une longueur de 30 mots ou moins.""",
        }
        key_findings = client.chat.completions.create(
            model=model_str,  # model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,  # set this to a low value, such as 0.2, to make the output more focused and deterministic. This will reduce the randomness in the generated response
            frequency_penalty=0.0,  # set this to a high positive value, such as 2.0, to penalize the model for repeating the same line verbatim. This will encourage the model to generate more diverse responses.
            presence_penalty=1.8,  # set this to a high positive value, such as 2.0, to penalize the model for introducing new tokens that do not appear in the input text. This will encourage the model to generate responses that are more closely related to the input text.
        )
        return key_findings

    def court_decision():
        question = {
            "en": f"""Classify the text as one of the following categories: Dismissed, Granted, Suspended, Rejected, or None. Simply mention the classification without going into details.""",
            "fr": "Classifiez le texte dans l'une des catégories suivantes : Rejeté, Accordé, Suspendu, Rejeté ou Aucun. Mentionnez simplement la classification sans entrer dans les détails.",
        }

        decision = client.chat.completions.create(
            model=model_str,  # model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,  # set this to a low value, such as 0.2, to make the output more focused and deterministic. This will reduce the randomness in the generated response
            frequency_penalty=0.0,  # set this to a high positive value, such as 2.0, to penalize the model for repeating the same line verbatim. This will encourage the model to generate more diverse responses.
            presence_penalty=1.8,  # set this to a high positive value, such as 2.0, to penalize the model for introducing new tokens that do not appear in the input text. This will encourage the model to generate responses that are more closely related to the input text.
        )
        return decision

    def find_keywords(findings):
        question = {
            "en": """Based on the findings, find multiple keywords of a total of around 5 to describe the case. Those keywords should have an emphasis on the type of application being challenged and what the applicant is seeking. Make sure your reply follows the format "keyword 1, keyword 2, ..., keyword 5". Make sure there's more than one keyword and don't include the court decision in the keyword. Do not reply anything else that doesn't follow this format""",
            "fr": """En fonction des résultats, trouvez plusieurs mots-clés d'un total d'environ 5 pour décrire le cas. Ces mots-clés devraient mettre l'accent sur le type d'application contestée et ce que le demandeur recherche. Assurez-vous que votre réponse suit le format "mot-clé 1, mot-clé 2, ..., mot-clé 5". Assurez-vous qu'il y ait plus d'un mot-clé et n'incluez pas la décision de la cour dans le mot-clé. Ne répondez rien d'autre qui ne suit pas ce format.""",
        }
        keywords = client.chat.completions.create(
            model=model_str,  # model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
                {"role": "user", "content": "### Findings: " + findings},
            ],
            temperature=0.1,  # set this to a low value, such as 0.2, to make the output more focused and deterministic. This will reduce the randomness in the generated response
            frequency_penalty=0.0,  # set this to a high positive value, such as 2.0, to penalize the model for repeating the same line verbatim. This will encourage the model to generate more diverse responses.
            presence_penalty=1.8,  # set this to a high positive value, such as 2.0, to penalize the model for introducing new tokens that do not appear in the input text. This will encourage the model to generate responses that are more closely related to the input text.
        )
        return keywords

    # CURRENTLY NOT IN USE
    def find_question_certification():
        question = {
            "en": f'Check to see if the applicant has proposed a question for certification. If yes, write down the question proposed and the outcome. If no, write down "The AI model was not able to find any proposed question for certification."',
            "fr": "Vérifiez si le demandeur a proposé une question pour la certification. Si oui, notez la question proposée et le résultat. Si non, notez \"Le modèle d'IA n'a pas pu trouver de question proposée pour la certification.\"",
        }

        question_cert = client.chat.completions.create(
            model=model_str,  # model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,  # set this to a low value, such as 0.2, to make the output more focused and deterministic. This will reduce the randomness in the generated response
            frequency_penalty=0.0,  # set this to a high positive value, such as 2.0, to penalize the model for repeating the same line verbatim. This will encourage the model to generate more diverse responses.
            presence_penalty=1.8,  # set this to a high positive value, such as 2.0, to penalize the model for introducing new tokens that do not appear in the input text. This will encourage the model to generate responses that are more closely related to the input text.
        )
        return question_cert

    # CURRENTLY NOT IN USE
    def find_test():
        question = {
            "en": 'If the application included any legal test, list out the full test and outcome of it. If no, simply write "The AI model was not able to find legal tests in this document".',
            "fr": "Si l'application incluait un test juridique, énumérez le test complet et son résultat. Sinon, écrivez simplement \"Le modèle d'IA n'a pas pu trouver de tests juridiques dans ce document\".",
        }

        test = client.chat.completions.create(
            model=model_str,  # model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,  # set this to a low value, such as 0.2, to make the output more focused and deterministic. This will reduce the randomness in the generated response
            frequency_penalty=0.0,  # set this to a high positive value, such as 2.0, to penalize the model for repeating the same line verbatim. This will encourage the model to generate more diverse responses.
            presence_penalty=1.8,  # set this to a high positive value, such as 2.0, to penalize the model for introducing new tokens that do not appear in the input text. This will encourage the model to generate responses that are more closely related to the input text.
        )
        return test

    # now run the functions in parallel
    with concurrent.futures.ThreadPoolExecutor() as executor:
        future1 = executor.submit(create_facts_and_court_sections)
        future2 = executor.submit(court_decision)
        # TODO decide wether or not we want to keep the 'Test' and 'Question for Certerfication' part of the report
        # future3 = executor.submit(find_question_certification)
        # future4 = executor.submit(find_test)

        facts = future1.result()
        decision = future2.result()
        # question_cert = future3.result()
        # test = future4.result()

    result_text = facts.choices[0].message.content
    sections = ["Facts", "Court"] if language == "en" else ["Faits", "Tribunal"]
    extract_sections_data = extract_sections(result_text, sections=sections)
    facts_text = extract_sections_data[sections[0]]
    findings_text = extract_sections_data[sections[1]]
    decision_text = decision.choices[0].message.content
    # question_cert_text = question_cert.choices[0].message.content
    # test_text = test.choices[0].message.content

    with concurrent.futures.ThreadPoolExecutor() as executor:
        future5 = executor.submit(create_key_findings, findings_text)
        future6 = executor.submit(find_keywords, findings_text)

        key_findings = future5.result()
        keywords = future6.result()

    key_findings_text = key_findings.choices[0].message.content
    keywords_text = keywords.choices[0].message.content

    tokens_consumed = (
        (facts.usage.total_tokens)
        + key_findings.usage.total_tokens
        + decision.usage.total_tokens
        + keywords.usage.total_tokens
        + extra_cost
    )
    estimated_cost = (
        math.ceil(tokens_consumed * settings.OPENAI_COST_PER_TOKEN * 10000) / 10000
    )

    context = {
        "date": date,
        "docket_number": docket_number.replace("\n", " ").strip(),
        "citation": citation,
        "applicant": applicant,
        "respondant": respondent,
        "title": title,
        "judge": judge,
        "estimated_cost": estimated_cost,
        "Facts": facts_text,
        "Court": findings_text,
        "Key_Findings": key_findings_text,
        "Keywords": keywords_text,
        "decision": decision_text,
        # "Question_Certification": question_cert_text,
        # "Test": test_text,
        "Court_Judgment": court_judgment,
    }

    # # pass each text value of context in the acronym replacement function
    for key in context:
        if isinstance(context[key], str):  # Ensure the value is a string
            context[key] = replace_with_acronyms(context[key], acronym_dict)

    # Define the template file and output file paths
    template_path = os.path.join(
        settings.BASE_DIR,
        "template_wizard",
        "templates",
        "template_wizard",
        "canlii_wizard",
        f"immigration_case_report_{language}.docx",
    )

    # Load the template file
    tpl = DocxTemplate(template_path)

    # # Render the template and save the output to the temporary file
    tpl.render(context)

    # Save the updated document
    file_content = io.BytesIO()
    tpl.save(file_content)

    # Move the "cursor" to the beginning of the buffer
    file_content.seek(0)

    return file_content, citation


def generate_atip_case_report(text, language):

    extra_cost = 0
    using_large_doc = False

    enc = tiktoken.get_encoding("cl100k_base")

    # Count how many characters are in the text
    token_count = len(enc.encode(text))

    # TODO: Change this to gpt-4 once we figure out how to deal with the gpt-4 delay errors
    token_count_max = 15000
    if token_count < token_count_max:
        content = text

    else:
        content = map_reduce(text, language, model_str)
        using_large_doc = True

    beginning_of_doc, end_of_doc = trim_to_tokens(text, max_tokens=token_count_max)

    # if we are using a large doc we'll just use the beginning of the doc to find what we want with the llm calls
    text_for_llm_meta_extraction = beginning_of_doc if using_large_doc else content
    decision_date, docket_number, citation, title, extra_cost_meta = extract_meta_info(
        text, text_for_llm_meta_extraction, model_str, language
    )
    extra_cost += extra_cost_meta

    # Extract judges name
    question = {
        "en": """Find all judges who worked on this case. This will usually be found in the sections "REASONS FOR JUDGMENT BY" and "CONCURRED IN BY". For your reply, simply output each name seperated by a comma. Don't write anything else.""",
        "fr": """Trouvez tous les juges qui ont travaillé sur cette affaire. Cela se trouve généralement dans les sections "JUGEMENT ET MOTIFS PUBLIC" ou "MOTIFS DU JUGEMENT", et "PRONONCÉS À L’AUDIENCE". Pour votre réponse, veuillez simplement afficher chaque nom séparé par une virgule. N'écrivez rien d'autre.""",
    }
    # use this end of the document if using a large document
    text_for_judge_extraction_llm_call = end_of_doc if using_large_doc else content
    extraction = client.chat.completions.create(
        model=model_str,
        messages=[
            {"role": "system", "content": system_prompt[language]},
            {
                "role": "user",
                "content": "### Content: " + text_for_judge_extraction_llm_call,
            },
            {"role": "user", "content": "### Question: " + question[language]},
        ],
        temperature=0.1,
        frequency_penalty=0.0,
        presence_penalty=1.8,
    )
    present = extraction.choices[0].message.content
    # extra_cost += extraction.usage.total_tokens

    def high_level_summary():
        question = {
            "en": "Using plain language, write a brief and high-level summary of the case in 1 or 2 sentences.",
            "fr": "En utilisant un langage simple, rédigez un résumé bref et de haut niveau du cas en 1 ou 2 phrases.",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content:" + content},
                {
                    "role": "user",
                    "content": "### Question:" + question[language],
                },
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )

        return response

    def summarize_legal_facts():
        question = {
            "en": "Summarize the legal facts in the case, as described in the provided text. Do not include the issues. I'll ask that later.",
            "fr": "Résumez les faits juridiques de l'affaire, tels qu'ils sont décrits dans le texte fourni. N'incluez pas les problèmes juridiques. Je les demanderai plus tard.",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content" + content},
                {
                    "role": "user",
                    "content": "### Question" + question[language],
                },
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )

        return response

    def legal_issues_list():
        question = {
            "en": "List the legal issues in the case, as described in the provided text. Do not include anything else. I'll ask that later.",
            "fr": "Énumérez les problèmes juridiques dans l'affaire, tels que décrits dans le texte fourni. Ne mentionnez rien d'autre. Je demanderai cela plus tard.",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content" + content},
                {
                    "role": "user",
                    "content": "### Question" + question[language],
                },
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )

        return response

    def decision_of_case():

        question = {
            "en": "Summarize the decision of the case in 60 words or less. Do not include the reasons for the decision. I'll ask that later.",
            "fr": "Résumez la décision de l'affaire en 60 mots ou moins. N'incluez pas les raisons de la décision. Je vous demanderai cela plus tard.",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content" + content},
                {
                    "role": "user",
                    "content": "### Question" + question[language],
                },
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )

        return response

    def reason_of_decision():

        question = {
            "en": "Summarize the reasons for the decision in the case. Do not include the decision itself. I'll ask that later.",
            "fr": "Résumez les raisons de la décision dans l'affaire. N'incluez pas la décision elle-même. Je la demanderai plus tard.",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content" + content},
                {
                    "role": "user",
                    "content": "### Question" + question[language],
                },
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )

        return response

    def find_PA():
        question = {
            "en": "Find the specific provisions of the Privacy Act or the PA discussed in the legal document. Provide the response in a list format. If none are found, just answer 'Not found' and don't write anything else.",
            "fr": """Trouvez les dispositions spécifiques de la Loi sur la protection des renseignements personnels ou de la LP mentionnées dans le document juridique. Fournissez la réponse sous forme de liste. Si aucune n'est trouvée, répondez simplement "Non trouvé" et n'écrivez rien d'autre.""",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content" + content},
                {
                    "role": "user",
                    "content": "### Question" + question[language],
                },
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )

        return response

    def find_ATIA():
        question = {
            "en": "Find the specific provisions of the Access to Information Act or the ATIA discussed in the provided legal document. Provide the response in a list format. If none are found, just answer 'Not found' and don't write anything else.",
            "fr": """Trouvez les dispositions spécifiques de la Loi sur l'accès à l'information ou la LAI discutées dans le document juridique fourni. Fournissez la réponse sous forme de liste. Si aucune n'est trouvée, répondez simplement "Non trouvé" et n'écrivez rien d'autre.""",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content:" + content},
                {
                    "role": "user",
                    "content": "### Question:" + question[language],
                },
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )

        return response

    # now run the functions in parallel
    with concurrent.futures.ThreadPoolExecutor() as executor:
        future1 = executor.submit(high_level_summary)
        future2 = executor.submit(summarize_legal_facts)
        future3 = executor.submit(legal_issues_list)
        future4 = executor.submit(decision_of_case)
        future5 = executor.submit(reason_of_decision)
        future6 = executor.submit(find_PA)
        future7 = executor.submit(find_ATIA)

        summary_response = future1.result()
        facts_response = future2.result()
        issues_response = future3.result()
        decision_response = future4.result()
        reasons_response = future5.result()
        privacy_provisions_response = future6.result()
        access_to_info_provisions_response = future7.result()

    summary = summary_response.choices[0].message.content
    facts = facts_response.choices[0].message.content
    issues = issues_response.choices[0].message.content
    decision = decision_response.choices[0].message.content
    reasons = reasons_response.choices[0].message.content
    privacy_provisions = privacy_provisions_response.choices[0].message.content
    access_to_info_provisions = access_to_info_provisions_response.choices[
        0
    ].message.content

    tokens_consumed = (
        (summary_response.usage.total_tokens)
        + (facts_response.usage.total_tokens)
        + (issues_response.usage.total_tokens)
        + (decision_response.usage.total_tokens)
        + (reasons_response.usage.total_tokens)
        + (privacy_provisions_response.usage.total_tokens)
        + (access_to_info_provisions_response.usage.total_tokens)
        + extra_cost
    )
    estimated_cost = (
        math.ceil((tokens_consumed) * settings.OPENAI_COST_PER_TOKEN * 10000) / 10000
    )

    context = {
        "title": title.replace("\n", " ").strip(),
        "citation": citation.replace("\n", " ").strip(),
        "decision_date": decision_date.replace("\n", " ").strip(),
        "docket_number": docket_number.replace("\n", " ").strip(),
        "privacy_provisions": markdown.markdown(privacy_provisions),
        "access_to_info_provisions": markdown.markdown(access_to_info_provisions),
        "summary": markdown.markdown(summary),
        "facts": markdown.markdown(facts),
        "issues": markdown.markdown(issues),
        "decision": markdown.markdown(decision),
        "reasons": markdown.markdown(reasons),
        "present": present.replace("\n", " ").strip(),
        "estimated_cost": estimated_cost,
    }

    # Define the template file and output file paths
    template_path = os.path.join(
        settings.BASE_DIR,
        "template_wizard",
        "templates",
        "template_wizard",
        "canlii_wizard",
        f"atip_case_report_{language}.docx",
    )

    # Load the template file
    tpl = DocxTemplate(template_path)

    # Make a copy of the context dictionary and prep the html content for docx insertion
    from html2docx import html2docx

    context_copy = context.copy()
    # context_copy["source_url"] = url
    context_copy["privacy_provisions"] = tpl.new_subdoc(
        html2docx(context["privacy_provisions"], title="")
    )
    context_copy["access_to_info_provisions"] = tpl.new_subdoc(
        html2docx(context["access_to_info_provisions"], title="")
    )
    context_copy["summary"] = tpl.new_subdoc(html2docx(context["summary"], title=""))
    context_copy["facts"] = tpl.new_subdoc(html2docx(context["facts"], title=""))
    context_copy["issues"] = tpl.new_subdoc(html2docx(context["issues"], title=""))
    context_copy["decision"] = tpl.new_subdoc(html2docx(context["decision"], title=""))
    context_copy["reasons"] = tpl.new_subdoc(html2docx(context["reasons"], title=""))

    # Render the template and save the output to a new file
    tpl.render(context_copy)
    file_content = io.BytesIO()
    tpl.save(file_content)

    # Move the "cursor" to the beginning of the buffer
    file_content.seek(0)

    return file_content, citation


def generate_general_case_report(text, language):

    extra_cost = 0
    using_large_doc = False

    enc = tiktoken.get_encoding("cl100k_base")

    # Count how many characters are in the text
    token_count = len(enc.encode(text))

    # TODO: Change this to gpt-4 once we figure out how to deal with the gpt-4 delay errors
    token_count_max = 15000
    if token_count < token_count_max:
        content = text

    else:
        content = map_reduce(text, language, model_str)
        using_large_doc = True

    beginning_of_doc, end_of_doc = trim_to_tokens(text, max_tokens=token_count_max)

    # if we are using a large doc we'll just use the beginning of the doc to find what we want with the llm calls
    text_for_llm_meta_extraction = beginning_of_doc if using_large_doc else content
    decision_date, docket_number, citation, title, extra_cost_meta = extract_meta_info(
        text, text_for_llm_meta_extraction, model_str, language
    )
    extra_cost += extra_cost_meta

    # function to find keywords
    def find_keywords():
        question = {
            "en": """Find multiple keywords of a total of around 5 to describe the case. Those keywords should have an emphasis on the type of application being challenged and what the applicant is seeking. Make sure your reply follows the format "keyword 1, keyword 2, ..., keyword 5". Make sure there's more than one keyword and don't include the court decision in the keyword. Do not reply anything else that doesn't follow this format""",
            "fr": """Trouvez plusieurs mots-clés d'un total d'environ 5 pour décrire le cas. Ces mots-clés devraient mettre l'accent sur le type d'application contestée et ce que le demandeur recherche. Assurez-vous que votre réponse suit le format "mot-clé 1, mot-clé 2, ..., mot-clé 5". Assurez-vous qu'il y ait plus d'un mot-clé et n'incluez pas la décision de la cour dans le mot-clé. Ne répondez rien d'autre qui ne suit pas ce format.""",
        }
        keywords = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,  # set this to a low value, such as 0.2, to make the output more focused and deterministic. This will reduce the randomness in the generated response
            frequency_penalty=0.0,  # set this to a high positive value, such as 2.0, to penalize the model for repeating the same line verbatim. This will encourage the model to generate more diverse responses.
            presence_penalty=1.8,  # set this to a high positive value, such as 2.0, to penalize the model for introducing new tokens that do not appear in the input text. This will encourage the model to generate responses that are more closely related to the input text.
        )
        return keywords

    # function to find cited legislation
    def find_cited_legislation():
        question = {
            "en": "Find all the legislation cited in the text. Provide the response in a list format. If none are found, just answer 'Not found' and don't write anything else.",
            "fr": "Trouvez toute la législation citée dans le texte. Fournissez la réponse sous forme d'une liste. Si aucune n'est trouvée, répondez simplement 'Non trouvé' et n'écrivez rien d'autre.",
        }
        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )
        return response

    # function to find facts and background of the case
    def find_facts_and_background():
        question = {
            "en": """Write down the facts and background of the case, and make sure to mention in your answers important sections and subsections of relevant acts or regulations found in the text. Note that the facts are the “who, when, what, where, and why” of the case. Hence, describe the history of the applicant, the events that led to the case, the legal claims, and defenses of each party. This section will all be written under 'Facts and Background:'. You will then write a second section which will present all key findings of the text. This second section will be under 'Key Findings:'.""",
            "fr": """Écrivez les faits et l'historique de l'affaire, et assurez-vous de mentionner dans vos réponses les sections et sous-sections importantes des lois ou règlements pertinents trouvés dans le texte. Notez que les faits sont le "qui, quand, quoi, où et pourquoi" de l'affaire. Décrivez donc l'historique du demandeur, les événements qui ont conduit à l'affaire, les revendications juridiques et les défenses de chaque partie. Cette section sera entièrement rédigée sous 'Faits:'. Vous écrirez ensuite une deuxième section qui présentera toutes les conclusions clés du texte. Cette deuxième section sera sous 'Conclusions clés:'.""",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )
        return response

    # function to generate summary of the case
    def generate_summary():
        question = {
            "en": "Write a brief summary of the case in 1 or 2 sentences.",
            "fr": "Rédigez un bref résumé de l'affaire en 1 ou 2 phrases.",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )
        return response

    # function to find legal issues of the case
    def find_legal_issues():
        question = {
            "en": "List the legal issues in the case.",
            "fr": "Énumérez les questions juridiques de l'affaire.",
        }

        response = client.chat.completions.create(
            model=model_str,
            messages=[
                {"role": "system", "content": system_prompt[language]},
                {"role": "user", "content": "### Content: " + content},
                {"role": "user", "content": "### Question: " + question[language]},
            ],
            temperature=0.1,
            frequency_penalty=0.0,
            presence_penalty=1.8,
        )
        return response

    # now run the functions in parallel
    with concurrent.futures.ThreadPoolExecutor() as executor:
        future1 = executor.submit(find_keywords)
        future2 = executor.submit(find_cited_legislation)
        future3 = executor.submit(find_facts_and_background)
        future4 = executor.submit(generate_summary)
        future5 = executor.submit(find_legal_issues)

        keywords = future1.result()
        cited_legislation = future2.result()
        facts_findings = future3.result()
        summary = future4.result()
        legal_issues = future5.result()

    keywords_text = keywords.choices[0].message.content
    cited_legislation_text = cited_legislation.choices[0].message.content
    summary_text = summary.choices[0].message.content
    legal_issues_text = legal_issues.choices[0].message.content
    facts_findings_text = facts_findings.choices[0].message.content

    sections = (
        ["Facts and Background", "Key Findings"]
        if language == "en"
        else ["Faits", "Conclusions clés"]
    )
    extract_sections_data = extract_sections(facts_findings_text, sections=sections)
    facts_background_text = extract_sections_data[sections[0]]
    key_findings_text = extract_sections_data[sections[1]]

    tokens_consumed = (
        (keywords.usage.total_tokens)
        + (cited_legislation.usage.total_tokens)
        + (facts_findings.usage.total_tokens)
        + (summary.usage.total_tokens)
        + (legal_issues.usage.total_tokens)
        + extra_cost
    )

    estimated_cost = (
        math.ceil((tokens_consumed) * settings.OPENAI_COST_PER_TOKEN * 10000) / 10000
    )

    context = {
        "title": title.replace("\n", " ").strip(),
        "citation": citation.replace("\n", " ").strip(),
        "decision_date": decision_date.replace("\n", " ").strip(),
        "docket_number": docket_number.replace("\n", " ").strip(),
        "keywords": keywords_text,
        "cited_legislation": cited_legislation_text,
        "facts_background": facts_background_text,
        "summary": summary_text,
        "legal_issues": legal_issues_text,
        "key_findings": key_findings_text,
        "estimated_cost": estimated_cost,
    }

    # Define the template file and output file paths
    template_path = os.path.join(
        settings.BASE_DIR,
        "template_wizard",
        "templates",
        "template_wizard",
        "canlii_wizard",
        f"general_report_{language}.docx",
    )

    # Load the template file
    tpl = DocxTemplate(template_path)

    # Make a copy of the context dictionary and prep the html content for docx insertion
    from html2docx import html2docx

    context_copy = context.copy()
    # context_copy["source_url"] = url
    context_copy["keywords"] = tpl.new_subdoc(html2docx(context["keywords"], title=""))
    context_copy["cited_legislation"] = tpl.new_subdoc(
        html2docx(context["cited_legislation"], title="")
    )
    context_copy["facts_background"] = tpl.new_subdoc(
        html2docx(context["facts_background"], title="")
    )
    context_copy["summary"] = tpl.new_subdoc(html2docx(context["summary"], title=""))
    context_copy["legal_issues"] = tpl.new_subdoc(
        html2docx(context["legal_issues"], title="")
    )
    context_copy["key_findings"] = tpl.new_subdoc(
        html2docx(context["key_findings"], title="")
    )

    # Render the template and save the output to a new file
    tpl.render(context_copy)
    file_content = io.BytesIO()
    tpl.save(file_content)

    # Move the "cursor" to the beginning of the buffer
    file_content.seek(0)

    return file_content, citation
