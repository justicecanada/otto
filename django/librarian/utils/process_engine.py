import csv
import hashlib
import io
import re
import subprocess
import tempfile
import uuid
from urllib.parse import urljoin

from django.conf import settings
from django.utils import timezone
from django.utils.translation import gettext as _

import filetype
import openpyxl  # Add this import for handling Excel files
import requests
import tiktoken
from bs4 import BeautifulSoup
from markdownify import markdownify
from structlog import get_logger

from librarian.utils.extract_emails import extract_msg
from librarian.utils.extract_zip import process_zip_file
from librarian.utils.markdown_splitter import MarkdownSplitter
from otto.models import Cost

logger = get_logger(__name__)


def markdownify_wrapper(text):
    """Wrapper to allow options to be passed to markdownify"""
    return markdownify(
        text,
        heading_style="ATX",
        bullets="*",
        strong_em_symbol="_",
        escape_misc=False,
    )


def fetch_from_url(url):
    try:
        r = requests.get(url, allow_redirects=True)
        content_type = guess_content_type(r.content, r.headers.get("content-type"), url)
        return r.content, content_type

    except Exception as e:
        logger.error(f"Failed to fetch from URL: {e}")
        raise Exception(f"Failed to fetch from URL: {e}")


def generate_hash(file_obj, block_size=65536):
    hasher = hashlib.sha256()
    # Always seek to the beginning before reading
    if hasattr(file_obj, "seek"):
        file_obj.seek(0)
    if hasattr(file_obj, "chunks"):
        for chunk in file_obj.chunks(block_size):
            hasher.update(chunk)
    else:
        while True:
            chunk = file_obj.read(block_size)
            if not chunk:
                break
            hasher.update(chunk)
    # Always seek back to the beginning after reading
    if hasattr(file_obj, "seek"):
        file_obj.seek(0)
    return hasher.hexdigest()


def extract_html_metadata(content):
    # Content is the binary data from response.content so convert it to a string
    soup = BeautifulSoup(decode_content(content), "html.parser")
    title_element = soup.find("title")
    title = title_element.get_text(strip=True) if title_element else None
    time_element = soup.find("time", {"property": "dateModified"})
    modified_at = (
        timezone.datetime.strptime(time_element.get_text(strip=True), "%Y-%m-%d")
        if time_element
        else None
    )
    return {
        "extracted_title": title,
        "extracted_modified_at": modified_at,
    }


def create_nodes(chunks, document):
    from llama_index.core.schema import NodeRelationship, RelatedNodeInfo, TextNode

    document_uuid = document.uuid_hex
    data_source_uuid = document.data_source.uuid_hex

    # Create a document (parent) node
    metadata = {"node_type": "document", "data_source_uuid": data_source_uuid}
    if document.title:
        metadata["title"] = document.title
    source = document.file_path or document.url or document.filename
    if source:
        metadata["source"] = source
    document_node = TextNode(text="", id_=document_uuid, metadata=metadata)
    document_node.relationships[NodeRelationship.SOURCE] = RelatedNodeInfo(
        node_id=document_node.node_id
    )

    # Create chunk (child) nodes
    metadata["node_type"] = "chunk"
    child_nodes = create_child_nodes(
        chunks,
        source_node_id=document_node.node_id,
        metadata=metadata,
    )

    # Update node properties
    new_nodes = [document_node] + child_nodes
    exclude_keys = ["page_range", "node_type", "data_source_uuid", "chunk_number"]
    for node in new_nodes:
        node.excluded_llm_metadata_keys = exclude_keys
        node.excluded_embed_metadata_keys = exclude_keys
        # The misspelling of "seperator" corresponds with the LlamaIndex codebase
        node.metadata_seperator = "\n"
        node.metadata_template = "{key}: {value}"
        node.text_template = "# {metadata_str}\ncontent:\n{content}\n\n"

    return new_nodes


def guess_content_type(
    content: str | bytes, content_type: str = "", path: str = ""
) -> str:
    # We consider these content types to be reliable and do not need further guessing
    trusted_content_types = [
        "application/pdf",
        "application/xml",
        "application/vnd.ms-outlook",
        "application/zip",
        "application/x-zip-compressed",
        "text/html",
        "text/markdown",
        "text/csv",
        "application/csv",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "officedocument.spreadsheetml.sheet",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "officedocument.presentationml.presentation",
        "image/jpeg",
        "image/jpg",
        "image/png",
        "image/bmp",
        "image/tiff",
        "image/tif",
        "image/heif",
        "image/heic",
    ]

    if content_type in trusted_content_types:
        return content_type

    if hasattr(content, "read"):
        content = content.read()

    if isinstance(content, bytes):
        # Explicitly handle Outlook emails
        if path.endswith(".msg"):
            return "application/vnd.ms-outlook"

        if path.endswith(".zip"):
            return "application/zip"

        if path.endswith(".docx"):
            return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        # Use filetype library to guess the content type
        kind = filetype.guess(content)
        if kind and not path.endswith(".md"):
            return kind.mime

        # Fallback to manual checks if filetype library fails
        try:
            content = content.decode("utf-8", errors="ignore")
        except UnicodeDecodeError:
            return content_type  # Unable to decode binary content

    if isinstance(content, str):
        if "text" in content_type and path.endswith(".md"):
            return "text/markdown"

        if content.startswith("<!DOCTYPE html>") or "<html" in content:
            return "text/html"

        if content.startswith("<?xml") or "<root" in content:
            return "application/xml"

        if content.startswith("{") or content.startswith("["):
            return "application/json"

    return content_type or "text/plain"


def get_process_engine_from_type(type):
    if "image" in type:
        return "IMAGE"
    elif "officedocument.wordprocessingml.document" in type:
        return "WORD"
    elif "officedocument.presentationml.presentation" in type:
        return "POWERPOINT"
    elif "application/vnd.ms-outlook" in type:
        return "OUTLOOK_MSG"
    elif "application/zip" in type or "application/x-zip-compressed" in type:
        return "ZIP"
    elif "application/pdf" in type:
        return "PDF"
    elif "text/html" in type:
        return "HTML"
    elif "text/markdown" in type:
        return "MARKDOWN"
    elif "text/csv" in type or "application/csv" in type:
        return "CSV"
    elif "spreadsheet" in type:
        return "EXCEL"
    else:
        return "TEXT"


def decode_content(
    content: bytes,
    encodings: list[str] = ["utf-8", "cp1252"],
) -> str:
    """
    Decode content with multiple encodings with fallback.

    Returns:
        Decoded string

    Raises:
        Exception: If content cannot be decoded with any of the provided encodings
    """
    for encoding in encodings:
        try:
            return content.decode(encoding)
        except UnicodeDecodeError as e:
            logger.debug(e)
            continue
    raise Exception(f"Failed to decode content with encodings: {encodings}")


def extract_markdown(
    content,
    process_engine,
    pdf_method="default",
    base_url=None,
    chunk_size=768,
    selector=None,
    root_document_id=None,
):
    try:
        enable_markdown = True
        if process_engine == "IMAGE":
            content = resize_to_azure_requirements(content)
            enable_markdown = False
            md = pdf_to_text_azure_read(content)
        elif process_engine == "PDF":
            if pdf_method == "default":
                enable_markdown = False
                md = pdf_to_text_pdfium(content)
                if len(md) < 10:
                    # Fallback to Azure Document Intelligence Read API to OCR
                    md = pdf_to_text_azure_read(content)
            elif pdf_method == "azure_layout":
                md = pdf_to_markdown_azure_layout(content)
            elif pdf_method == "azure_read":
                enable_markdown = False
                md = pdf_to_text_azure_read(content)
        elif process_engine == "WORD":
            md = docx_to_markdown(content)
        elif process_engine == "POWERPOINT":
            md = pptx_to_markdown(content)
        elif process_engine == "HTML":
            md = html_to_markdown(decode_content(content), base_url, selector)
        elif process_engine == "MARKDOWN":
            md = decode_content(content)
        elif process_engine == "OUTLOOK_MSG":
            enable_markdown = False
            md = extract_msg(content, root_document_id)
        elif process_engine == "ZIP":
            enable_markdown = False
            md = process_zip_file(content, root_document_id)
        elif process_engine == "CSV":
            md = csv_to_markdown(content)
        elif process_engine == "EXCEL":
            md = excel_to_markdown(content)
        else:
            enable_markdown = False
            try:
                md = decode_content(content)
            except Exception as e:
                raise e

        md = remove_nul_characters(md)

        # Strip leading/trailing whitespace; replace all >2 line breaks with 2 line breaks
        md = re.sub(r"\n{3,}", "\n\n", md.strip())

        # Divide the markdown into chunks
        try:
            md_splitter = MarkdownSplitter(
                chunk_size=chunk_size, chunk_overlap=0, enable_markdown=enable_markdown
            )
            md_chunks = md_splitter.split_markdown(md)
        except Exception as e:
            logger.debug("Error splitting markdown using MarkdownSplitter:")
            logger.error(e)
            # Fallback to simpler method
            from llama_index.core.node_parser import SentenceSplitter

            sentence_splitter = SentenceSplitter(
                chunk_size=chunk_size, chunk_overlap=min(chunk_size // 4, 100)
            )
            md_chunks = sentence_splitter.split_text(md)
        return md, md_chunks
    except Exception as e:
        logger.error(f"Error in extract_markdown: {str(e)}")
        raise


def pdf_to_markdown_azure_layout(content):
    html = _pdf_to_html_azure_layout(content)
    return _convert_html_to_markdown(html)


def pdf_to_text_pdfium(content):
    # Fast and cheap, but no OCR or layout analysis
    import pypdfium2 as pdfium

    from otto.utils.common import pdfium_lock

    with pdfium_lock:
        try:
            pdf = pdfium.PdfDocument(content)
        except Exception as e:
            logger.error(f"Failed to extract text from PDF file: {e}")
            raise Exception(_("Corrupt PDF file."))

        text = ""
        for i, page in enumerate(pdf):
            text_page = page.get_textpage()
            text_content = text_page.get_text_range()
            if text_content:
                text += f"<page_{i+1}>\n"
                text += text_content + "\n"
                text += f"</page_{i+1}>\n"
            # PyPDFium does not cleanup its resources automatically. Ensures memory freed.
            text_page.close()
        pdf.close()

    return text


def html_to_markdown(content, base_url=None, selector=None):
    return _convert_html_to_markdown(content, base_url, selector)


def remove_nul_characters(text):
    """Remove NUL (0x00) characters from the text."""
    return text.replace("\x00", "")


def msg_to_markdown(content):
    with tempfile.NamedTemporaryFile(suffix=".msg") as temp_file:
        temp_file.write(content)
        temp_file_path = temp_file.name
        try:
            md = subprocess.check_output(
                ["python", "-m", "extract_msg", "--dump-stdout", temp_file_path]
            ).decode("utf-8")
        except subprocess.CalledProcessError as e:
            logger.error(f"Command failed with exit code {e.returncode}")
            logger.error(f"Output: {e.output.decode('utf-8')}")
            md = ""
        except Exception as e:
            logger.error(f"Failed to extract text from Outlook email: {e}")
            md = ""
        return md


def docx_to_markdown(content):

    import mammoth

    with io.BytesIO(content) as docx_file:
        try:
            result = mammoth.convert_to_html(docx_file)
        except Exception as e:
            logger.error(f"Failed to extract text from .docx file: {e}")
            raise Exception(_("Corrupt docx file."))
    html = result.value

    return _convert_html_to_markdown(html)


def pptx_to_markdown(content):
    import pptx

    with io.BytesIO(content) as ppt_file:
        try:
            prs = pptx.Presentation(ppt_file)
        except Exception as e:
            logger.error(f"Failed to extract text from .pptx file: {e}")
            raise Exception(_("Corrupt pptx file."))

    # extract text from each slide
    all_html = ""
    for i, slide in enumerate(prs.slides):
        html = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                html += "<p>"
                for run in paragraph.runs:
                    html += run.text
                html += "</p>"
        if len(slide.notes_slide.notes_text_frame.paragraphs) > 0:
            html += f"<h6>Presenter notes:</h6>"
            for note in slide.notes_slide.notes_text_frame.paragraphs:
                html += "<p>"
                for run in note.runs:
                    html += run.text
                html += "</p>"
        if html:
            all_html += f"<page_{i+1}>\n{html}\n</page_{i+1}>\n"

    return _convert_html_to_markdown(all_html)


def create_child_nodes(chunks, source_node_id, metadata=None):
    from llama_index.core.schema import NodeRelationship, RelatedNodeInfo, TextNode

    nodes = []
    for i, text in enumerate(chunks):

        node = TextNode(text=text, id_=str(uuid.uuid4()))

        node.metadata = dict(metadata, chunk_number=i)
        node.relationships[NodeRelationship.SOURCE] = RelatedNodeInfo(
            node_id=source_node_id
        )
        nodes.append(node)

    # Handle the case when there's only one or zero elements
    if len(chunks) < 2:
        return nodes

    # Set relationships
    for i in range(len(nodes) - 1):
        nodes[i].relationships[NodeRelationship.NEXT] = RelatedNodeInfo(
            node_id=nodes[i + 1].node_id
        )
        nodes[i + 1].relationships[NodeRelationship.PREVIOUS] = RelatedNodeInfo(
            node_id=nodes[i].node_id
        )

    return nodes


def token_count(string: str, model: str = "gpt-4") -> int:
    """Returns the number of tokens in a text string."""
    encoding = tiktoken.encoding_for_model(model)
    num_tokens = len(encoding.encode(string))
    return num_tokens


def _remove_ignored_tags(text):
    # remove any javascript, css, images, svg, and comments from self.text
    text = re.sub(r"<script.*?</script>", "", text, flags=re.DOTALL)
    text = re.sub(r"<style.*?</style>", "", text, flags=re.DOTALL)
    text = re.sub(r"<!--.*?-->", "", text, flags=re.DOTALL)
    text = re.sub(r"<img.*?>", "", text, flags=re.DOTALL)
    text = re.sub(r"<svg.*?</svg>", "", text, flags=re.DOTALL)
    # remove any attribute tags that start with javascript:
    text = re.sub(r"<[^>]+javascript:.*?>", "", text, flags=re.DOTALL)
    # remove any empty html tags from self.text
    text = re.sub(r"<[^/>][^>]*>\s*</[^>]+>", "", text, flags=re.DOTALL)

    # remove any header/footer/nav tags and the content within them
    text = re.sub(r"<header.*?</header>", "", text, flags=re.DOTALL)
    text = re.sub(r"<footer.*?</footer>", "", text, flags=re.DOTALL)
    text = re.sub(r"<nav.*?</nav>", "", text, flags=re.DOTALL)

    # Remove all the line breaks, carriage returns, and tabs
    text = re.sub(r"[\n\r\t]", "", text)

    return text


def _convert_html_to_markdown(
    source_html: str, base_url: str = None, selector: str = None
) -> str:
    """Converts HTML to markdown, preserving <page_x> tags in the markdown output."""
    page_open_tags = re.findall(r"<page_\d+>", source_html)
    # When page tags (e.g. "<page_1">) are present, run this step separately for each
    # of the page contents and combine the results
    if page_open_tags:
        combined_md = ""
        for opening_tag in page_open_tags:
            closing_tag = opening_tag.replace("<", "</")
            page_html_contents = re.search(
                f"{opening_tag}(.*){closing_tag}", source_html, re.DOTALL
            ).group(1)
            page_md = _convert_html_to_markdown(page_html_contents, base_url)
            combined_md += f"{opening_tag}\n{page_md}\n{closing_tag}\n"
        return combined_md

    soup = BeautifulSoup(source_html, "html.parser")
    if soup.find("body"):
        soup = soup.find("body")

    if selector:
        selected_html = soup.select_one(selector)
        if selected_html:
            soup = selected_html
        else:
            logger.warning(f"Selector {selector} not found in HTML")

    if not base_url:
        # find all anchor tags
        for anchor in soup.find_all("a"):
            # get the href attribute value
            href = anchor.get("href")
            # convert relative URLs to absolute URLs
            if href and not href.startswith("http"):
                absolute_url = urljoin(base_url, href)
                anchor["href"] = absolute_url

    # Replace <caption> elements with <h6> so that they get capture in breadcrumbs
    for caption in soup.find_all("caption"):
        caption.name = "h6"

    text = _remove_ignored_tags(str(soup))

    markdown = markdownify_wrapper(text).strip()
    return markdown


def _pdf_to_html_azure_layout(content):
    from azure.ai.formrecognizer import DocumentAnalysisClient
    from azure.core.credentials import AzureKeyCredential
    from shapely.geometry import Polygon

    # Note: This method handles scanned PDFs, images, and handwritten text but is $$$

    document_analysis_client = DocumentAnalysisClient(
        endpoint=settings.AZURE_COGNITIVE_SERVICE_ENDPOINT,
        credential=AzureKeyCredential(settings.AZURE_COGNITIVE_SERVICE_KEY),
    )

    poller = document_analysis_client.begin_analyze_document("prebuilt-layout", content)
    result = poller.result()

    num_pages = len(result.pages)
    cost = Cost.objects.new(cost_type="doc-ai-prebuilt", count=num_pages)

    # Extract table bounding regions
    table_bounding_regions = []
    for table in result.tables:
        for cell in table.cells:
            table_bounding_regions.append(cell.bounding_regions[0])

    table_chunks = []
    for _, table in enumerate(result.tables):
        page_number = table.bounding_regions[0].page_number

        # Generate table HTML syntax
        table_html = "<table>"
        for row in range(table.row_count):
            table_html += "<tr>"
            for col in range(table.column_count):
                try:
                    cell = next(
                        cell
                        for cell in table.cells
                        if cell.row_index == row and cell.column_index == col
                    )
                    table_html += "<td>{}</td>".format(cell.content)
                except StopIteration:
                    table_html += "<td></td>"
            table_html += "</tr>"
        table_html += "</table>"

        chunk = {
            "page_number": page_number,
            "x": table.bounding_regions[0].polygon[0].x,
            "y": table.bounding_regions[0].polygon[0].y,
            "text": table_html,
        }
        table_chunks.append(chunk)

    p_chunks = []
    for paragraph in result.paragraphs:
        paragraph_page_number = paragraph.bounding_regions[0].page_number
        paragraph_polygon = Polygon(
            [(point.x, point.y) for point in paragraph.bounding_regions[0].polygon]
        )

        # Check intersection between paragraph and table cells
        if any(
            paragraph_polygon.intersects(
                Polygon([point.x, point.y] for point in cell.polygon)
            )
            for cell in table_bounding_regions
            if cell.page_number == paragraph_page_number
        ):
            continue

        # If text contains words like :selected:, :checked:, or :unchecked:, then skip it
        if any(
            word in paragraph.content
            for word in [":selected:", ":checked:", ":unchecked:"]
        ):
            continue

        # Create Chunk object and append to chunks list
        chunk = {
            "page_number": paragraph_page_number,
            "x": paragraph_polygon.bounds[0],
            "y": paragraph_polygon.bounds[1],
            "text": "<p>" + paragraph.content + "</p>",
        }
        p_chunks.append(chunk)

    chunks = table_chunks + p_chunks

    # Sort chunks by page number, then by y coordinate, then by x coordinate
    chunks = sorted(
        chunks, key=lambda item: (item.get("page_number"), item.get("y"), item.get("x"))
    )
    html = ""
    cur_page = None
    for _, chunk in enumerate(chunks, 1):
        page_start_tag = f"\n<page_{chunk.get('page_number')}>\n"
        page_end_tag = f"\n</page_{chunk.get('page_number')}>\n"
        prev_end_tag = f"\n</page_{cur_page}>\n" if cur_page is not None else ""
        if chunk.get("page_number") != cur_page:
            if cur_page is not None:
                html += prev_end_tag
            cur_page = chunk.get("page_number")
            html += page_start_tag
        html += chunk.get("text")

    if cur_page is not None and chunks:
        html += page_end_tag

    return html


def pdf_to_text_azure_read(content: bytes) -> str:

    from azure.ai.formrecognizer import DocumentAnalysisClient
    from azure.core.credentials import AzureKeyCredential

    document_analysis_client = DocumentAnalysisClient(
        endpoint=settings.AZURE_COGNITIVE_SERVICE_ENDPOINT,
        credential=AzureKeyCredential(settings.AZURE_COGNITIVE_SERVICE_KEY),
    )

    poller = document_analysis_client.begin_analyze_document("prebuilt-read", content)
    result = poller.result()

    num_pages = len(result.pages)
    cost = Cost.objects.new(cost_type="doc-ai-read", count=num_pages)

    p_chunks = []
    for page in result.pages:
        for line in page.lines:
            chunk = {
                "page_number": page.page_number,
                "text": line.content + "\n",
            }
            p_chunks.append(chunk)

    text = ""
    cur_page = None
    for _, chunk in enumerate(p_chunks, 1):
        page_start_tag = f"\n<page_{chunk.get('page_number')}>\n"
        page_end_tag = f"\n</page_{chunk.get('page_number')}>\n"
        prev_end_tag = f"\n</page_{cur_page}>\n" if cur_page is not None else ""
        if chunk.get("page_number") != cur_page:
            if cur_page is not None:
                text = text.strip() + prev_end_tag
            cur_page = chunk.get("page_number")
            text = text.strip() + page_start_tag
        text += chunk.get("text")

    if cur_page is not None and p_chunks:
        text = text.strip() + page_end_tag

    return text


def csv_to_markdown(content):
    """Convert CSV content to markdown table."""
    try:
        with io.StringIO(content.decode("utf-8")) as csv_file:
            reader = csv.reader(csv_file)
            rows = list(reader)
    except Exception as e:
        logger.error(f"Failed to extract text from CSV file: {e}")
        raise Exception(_("Corrupt CSV file."))

    if not rows:
        return ""

    header = rows[0]
    table = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * len(header)) + " |",
    ]
    for row in rows[1:]:
        table.append("| " + " | ".join(row) + " |")

    return "\n".join(table)


def excel_to_markdown(content):
    """Convert Excel content to markdown tables."""
    try:
        workbook = openpyxl.load_workbook(io.BytesIO(content))
    except Exception as e:
        logger.error(f"Failed to extract text from Excel file: {e}")
        raise Exception(_("Corrupt Excel file."))

    markdown = ""
    for sheet in workbook.sheetnames:
        markdown += f"# {sheet}\n\n"
        sheet_obj = workbook[sheet]
        rows = list(sheet_obj.values)
        if not rows:
            continue
        header = rows[0]
        table = [
            "| " + " | ".join(map(str, header)) + " |",
            "| " + " | ".join(["---"] * len(header)) + " |",
        ]
        for row in rows[1:]:
            table.append("| " + " | ".join(map(str, row)) + " |")
        markdown += "\n".join(table) + "\n\n"
    return markdown


def resize_to_azure_requirements(content):
    from PIL import Image

    if isinstance(content, Image.Image):
        image = content
    else:
        with io.BytesIO(content) as image_file:
            image = Image.open(image_file)
            image.load()

    width, height = image.size
    if width < 50 or height < 50:
        # Resize to at least 50 pixels
        if width <= height:
            new_width = 50
            new_height = int(height * (50 / width))
        else:
            new_height = 50
            new_width = int(width * (50 / height))
    elif width > 10000 or height > 10000:
        # Resize to max 10000 pixels
        if width >= height:
            new_width = 10000
            new_height = int(height * (10000 / width))
        else:
            new_height = 10000
            new_width = int(width * (10000 / height))
    else:
        if isinstance(content, Image.Image):
            new_width, new_height = width, height
        else:
            return content
    # Edge case: insanely wide or tall images. Don't maintain proportions.
    new_width = min(new_width, 10000)
    new_height = min(new_height, 10000)
    new_width = max(new_width, 50)
    new_height = max(new_height, 50)
    image = image.resize((new_width, new_height))

    if isinstance(content, Image.Image):
        return image
    else:
        with io.BytesIO() as output:
            image.save(output, format="PNG")
            content = output.getvalue()
            return content
