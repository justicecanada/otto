import io
import os
import shutil
import uuid
from datetime import datetime
from unittest.mock import MagicMock

from django.conf import settings
from django.contrib.auth.models import Group
from django.core.files.base import ContentFile
from django.core.management import call_command
from django.test import override_settings

import pytest
import pytest_asyncio
from asgiref.sync import sync_to_async
from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from reportlab.pdfgen import canvas

from text_extractor.models import OutputFile

pytest_plugins = ("pytest_asyncio",)

this_dir = os.path.dirname(os.path.abspath(__file__))


@pytest.fixture(scope="function", autouse=True)
def set_test_media():
    # Define the test media directory
    test_media_dir = os.path.join(settings.BASE_DIR, "test_media")

    storages = settings.STORAGES.copy()
    storages["default"]["LOCATION"] = test_media_dir

    # Ensure the test media directory is clean
    if os.path.exists(test_media_dir):
        shutil.rmtree(test_media_dir)
    os.makedirs(test_media_dir)

    # Use override_settings to set MEDIA_ROOT
    with override_settings(STORAGES=storages, MEDIA_ROOT=test_media_dir):
        yield  # This allows the tests to run

    # Cleanup after tests
    shutil.rmtree(test_media_dir)


@pytest_asyncio.fixture(scope="session")
async def django_db_setup(django_db_setup, django_db_blocker):
    def _inner():
        with django_db_blocker.unblock():
            call_command(
                "reset_app_data",
                "groups",
                "apps",
                "security_labels",
                "library_mini",
                "cost_types",
                "presets",
            )
            # Process the Wikipedia document only
            from chat.llm import OttoLLM
            from librarian.models import Document
            from librarian.tasks import process_document_helper

            test_document = Document.objects.get(
                url="https://en.wikipedia.org/wiki/Glyph"
            )
            process_document_helper(test_document, OttoLLM())

    return await sync_to_async(_inner)()


@pytest.fixture()
def load_example_pdf(django_db_blocker):
    from chat.llm import OttoLLM
    from librarian.models import DataSource, Document, SavedFile
    from librarian.tasks import process_document_helper

    with open(os.path.join(this_dir, "librarian/test_files/example.pdf"), "rb") as f:
        with django_db_blocker.unblock():
            pdf_file = ContentFile(f.read(), name="example.pdf")
            saved_file = SavedFile.objects.create(file=pdf_file)
            d = Document.objects.create(
                saved_file=saved_file,
                filename="example.pdf",
                data_source=DataSource.objects.get(name_en="Wikipedia"),
            )
            process_document_helper(d, OttoLLM())


@pytest.fixture()
def all_apps_user(db, django_user_model):
    def new_user(username="all_apps_user"):
        user = django_user_model.objects.create_user(
            upn=f"{username}.lastname@example.com",
            oid=f"{username}_oid",
            email=f"{username}@example.com",
        )
        user.groups.add(Group.objects.get(name="Otto admin"))
        # Accept the terms
        user.accepted_terms_date = datetime.now()
        user.save()
        return user

    return new_user


@pytest.fixture()
def basic_user(db, django_user_model):
    def new_user(username="basic_user", accept_terms=False):
        user = django_user_model.objects.create_user(
            upn=f"{username}.lastname@example.com",
            oid=f"{username}_oid",
            email=f"{username}@example.com",
            accepted_terms_date=datetime.now() if accept_terms else None,
        )
        return user

    return new_user


@pytest.fixture
def mock_pdf_file():
    filename = "temp_file1.pdf"
    c = canvas.Canvas(filename)
    for i in range(3):  # Create 3 pages
        c.drawString(100, 100, f"Page {i+1}")
        c.showPage()
    c.save()

    with open(filename, "rb") as f:
        yield f
    os.remove(filename)


@pytest.fixture
def mock_pdf_file2():
    filename = "temp_file2.pdf"
    c = canvas.Canvas(filename)
    for i in range(10):  # Create 10 pages
        c.drawString(100, 100, f"Page {i+1}")
        c.showPage()
    c.save()

    with open(filename, "rb") as f:
        yield f
    os.remove(filename)


# yields filename and content
@pytest.fixture
def mock_pdf_file3():
    filename = "temp_file1.pdf"
    c = canvas.Canvas(filename)
    for i in range(3):  # Create 3 pages
        c.drawString(100, 100, f"Page {i+1}")
        c.showPage()
    c.save()

    with open(filename, "rb") as f:
        content = f.read()
        yield filename, content
    os.remove(filename)


@pytest.fixture
def mock_image_file(filename="temp_image.jpg"):
    mock_file = MagicMock()
    mock_file.name = filename
    yield mock_file


@pytest.fixture
def mock_image_file2():
    img = Image.new("RGB", (1000, 500), "white")
    return img


# yields filename and content
@pytest.fixture
def mock_image_file3():
    filename = "temp_image.jpg"
    # Create a simple image
    image = Image.new("RGB", (100, 100), color="red")
    draw = ImageDraw.Draw(image)

    # Draw the letter "R" in black color
    font = ImageFont.load_default()

    draw.text((25, 25), "RIF drawing", fill="black", font=font)
    image.save(filename)

    # Open the file in binary read mode and return the file object and its content
    with open(filename, "rb") as f:
        content = f.read()
        yield filename, content
    os.remove(filename)


@pytest.fixture
def mock_unsupported_file():
    mock_file = MagicMock()
    mock_file.name = "temp_unsupported.txt"
    yield mock_file


# Mocking file objects with a .name attribute
class MockFile:
    def __init__(self, name, total_page_num):
        self.name = name


@pytest.fixture
def process_ocr_document_mock(mocker):
    # Mock the Celery task's delay method
    mock_delay = mocker.patch("text_extractor.views.process_ocr_document.delay")
    mock_task = MagicMock()
    mock_task.id = "mock_task_id"
    mock_delay.return_value = mock_task

    # Mock the AsyncResult
    mock_async_result = mocker.patch(
        "text_extractor.views.process_ocr_document.AsyncResult"
    )
    mock_result_instance = MagicMock()
    # Set the return value of result.get()
    mock_result_instance.get.return_value = (
        b"pdf_bytes_content",
        "txt_file_content",
        0.05,
        "input_name",
    )
    mock_async_result.return_value = mock_result_instance

    return mock_delay, mock_async_result


@pytest.fixture
def content_file_mock(mocker):
    original_content_file = ContentFile
    mock_content_file = mocker.patch(
        "django.core.files.base.ContentFile", side_effect=original_content_file
    )
    return mock_content_file


@pytest.fixture
def basic_feedback():
    from django.utils import timezone

    from otto.forms import FeedbackForm
    from otto.models import Feedback

    def new_feedback_form(user):
        date_and_time = timezone.now().strftime("%Y%m%d-%H%M%S")
        feedback = Feedback(
            feedback_type=Feedback.FEEDBACK_TYPE_CHOICES[0][0],
            feedback_message="Test Message",
            app="Otto",
            modified_by=user,
            created_by=user,
            created_at=date_and_time,
            modified_on=date_and_time,
            otto_version="v0",
        )
        return feedback

    return new_feedback_form


@pytest.fixture
def output_file():
    pdf_mock = MagicMock()
    pdf_mock.name = "test.pdf"
    pdf_mock.open.return_value.__enter__.return_value.read.return_value = b"PDF content"

    txt_mock = MagicMock()
    txt_mock.name = "test.txt"
    txt_mock.open.return_value.__enter__.return_value.read.return_value = b"TXT content"

    output_file = MagicMock(spec=OutputFile)
    output_file.celery_task_ids = [str(uuid.uuid4())]
    output_file.pdf_file = pdf_mock
    output_file.txt_file = txt_mock
    output_file.file_name = "test_document"
    output_file.usd_cost = 0

    return output_file


@pytest.fixture
def sample_docx():
    doc = Document()
    doc.add_heading("Test Heading", 0)
    doc.add_paragraph("Test paragraph")
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


@pytest.fixture
def sample_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "Test Content"
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


@pytest.fixture
def sample_excel():
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Header1"
    ws["B1"] = "Header2"
    ws["A2"] = "Value1"
    ws["B2"] = "Value2"
    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


@pytest.fixture
def sample_csv():
    return b"Header1,Header2\nValue1,Value2"


@pytest.fixture
def sample_pdf():
    # Mock PDF content for testing
    return b"%PDF-1.4\n..."
